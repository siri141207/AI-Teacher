import os
from pypdf import PdfReader
import docx
from pptx import Presentation


def extract_pdf_text(file_path: str) -> str:
    reader = PdfReader(file_path)
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages.append(f"[Page {i+1}]\n{text.strip()}")
    return "\n\n".join(pages)


def extract_docx_text(file_path: str) -> str:
    doc = docx.Document(file_path)
    paragraphs = []
    for p in doc.paragraphs:
        if p.text.strip():
            paragraphs.append(p.text.strip())
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return "\n\n".join(paragraphs)


def extract_pptx_text(file_path: str) -> str:
    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        slide_parts.append(line)
        if slide_parts:
            slides_text.append(f"[Slide {i+1}]\n" + "\n".join(slide_parts))
    return "\n\n".join(slides_text)


def extract_plain_text(file_path: str) -> str:
    encodings = ["utf-8", "latin-1", "cp1252"]
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    with open(file_path, "r", errors="ignore") as f:
        return f.read()


def extract_document_text(file_path: str) -> str:
    """Universal extractor for PDF, DOCX, PPTX, TXT, MD documents."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_pdf_text(file_path)
    elif ext in [".docx", ".doc"]:
        return extract_docx_text(file_path)
    elif ext in [".pptx", ".ppt"]:
        return extract_pptx_text(file_path)
    elif ext in [".txt", ".md", ".csv", ".json", ".rtf"]:
        return extract_plain_text(file_path)
    else:
        # Fallback to plain text reader
        return extract_plain_text(file_path)